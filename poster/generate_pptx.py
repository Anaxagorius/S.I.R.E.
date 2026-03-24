#!/usr/bin/env python3
"""
S.I.R.E. PowerPoint Poster Generator
Combines the poster content and styling with the layoutA_fixed_structure.pptx
3-panel template.

Run with:  python3 poster/generate_pptx.py
Output:    poster/SIRE_Poster_Presentation.pptx
"""

import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ─────────────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT  = os.path.dirname(SCRIPT_DIR)
TEMPLATE   = os.path.join(REPO_ROOT, "layoutA_fixed_structure.pptx")
OUTPUT     = os.path.join(SCRIPT_DIR, "SIRE_Poster_Presentation.pptx")

# ── Colour palette (matches poster / SIRE brand) ──────────────────────────────
RED        = RGBColor(0xE6, 0x39, 0x46)
RED_DARK   = RGBColor(0xC1, 0x12, 0x1F)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1A, 0x1A, 0x1A)
MUTED      = RGBColor(0x55, 0x55, 0x55)
BORDER     = RGBColor(0xDD, 0xDD, 0xDD)
LEFT_BG    = RGBColor(0xFF, 0xFA, 0xFA)
RIGHT_BG   = RGBColor(0xFA, 0xFA, 0xFA)

SWOT_S_BG  = RGBColor(0xE8, 0xF8, 0xEF); SWOT_S_BD = RGBColor(0xA3, 0xD9, 0xB8)
SWOT_S_TXT = RGBColor(0x1A, 0x5C, 0x36)
SWOT_W_BG  = RGBColor(0xFF, 0xF0, 0xE0); SWOT_W_BD = RGBColor(0xF4, 0xC4, 0x8A)
SWOT_W_TXT = RGBColor(0x6B, 0x38, 0x00)
SWOT_O_BG  = RGBColor(0xE8, 0xF0, 0xFF); SWOT_O_BD = RGBColor(0xA3, 0xB8, 0xF4)
SWOT_O_TXT = RGBColor(0x1A, 0x2E, 0x6B)
SWOT_T_BG  = RGBColor(0xFF, 0xF0, 0xF0); SWOT_T_BD = RGBColor(0xF4, 0xA3, 0xA3)
SWOT_T_TXT = RGBColor(0x6B, 0x1A, 0x1A)

AB_BLUE_BG = RGBColor(0xDA, 0xEE, 0xFF); AB_BLUE_BD = RGBColor(0x5B, 0xA8, 0xD4)
AB_BLUE_TX = RGBColor(0x0A, 0x3C, 0x5C)
AB_GREEN_BG= RGBColor(0xDA, 0xFF, 0xEF); AB_GREEN_BD= RGBColor(0x38, 0xB8, 0x85)
AB_GREEN_TX= RGBColor(0x0A, 0x40, 0x28)
AB_RED_BG  = RGBColor(0xFF, 0xE8, 0xE8); AB_RED_BD  = RGBColor(0xE0, 0x60, 0x60)
AB_RED_TX  = RGBColor(0x6B, 0x00, 0x00)
AB_GOLD_BG = RGBColor(0xFF, 0xF5, 0xD6); AB_GOLD_BD = RGBColor(0xD4, 0xA0, 0x17)
AB_GOLD_TX = RGBColor(0x5C, 0x3A, 0x00)
AB_PURP_BG = RGBColor(0xF2, 0xE8, 0xFF); AB_PURP_BD = RGBColor(0xA0, 0x60, 0xD4)
AB_PURP_TX = RGBColor(0x3A, 0x00, 0x5C)

FE_BG  = RGBColor(0xDA, 0xEE, 0xFF); FE_TX  = RGBColor(0x0A, 0x3C, 0x5C)
BE_BG  = RGBColor(0xDA, 0xFF, 0xEF); BE_TX  = RGBColor(0x0A, 0x40, 0x28)
RT_BG  = RGBColor(0xFF, 0xF5, 0xD6); RT_TX  = RGBColor(0x5C, 0x3A, 0x00)
OPS_BG = RGBColor(0xFF, 0xE8, 0xE8); OPS_TX = RGBColor(0x6B, 0x00, 0x00)

PINK_BG   = RGBColor(0xFF, 0xF0, 0xF0)
BLUE_M_BG = RGBColor(0xE8, 0xF0, 0xFF)
GRN_M_BG  = RGBColor(0xE8, 0xF8, 0xEF)
YEL_M_BG  = RGBColor(0xFF, 0xFD, 0xE8)
MET_RED   = RGBColor(0xE6, 0x39, 0x46)
MET_BLUE  = RGBColor(0x4A, 0x90, 0xD9)
MET_GREEN = RGBColor(0x06, 0xB8, 0x8A)
MET_GOLD  = RGBColor(0xD4, 0xA0, 0x17)

AV_RED  = RGBColor(0xE6, 0x39, 0x46)
AV_BLUE = RGBColor(0x4A, 0x90, 0xD9)
AV_TEAL = RGBColor(0x06, 0xB8, 0x8A)
AV_GOLD = RGBColor(0xD4, 0xA0, 0x17)

# ── EMU helpers ───────────────────────────────────────────────────────────────
def _i(inches):  return Inches(inches)
def _p(points):  return Pt(points)


# ─────────────────────────────────────────────────────────────────────────────
# Low-level shape helpers
# ─────────────────────────────────────────────────────────────────────────────

def _set_solid_fill(shape, rgb: RGBColor):
    """Force a solid fill on any shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _set_line(shape, rgb: RGBColor, width_pt: float = 1.0):
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)


def add_rect(slide, l, t, w, h, bg: RGBColor = None,
             border: RGBColor = None, border_pt: float = 1.0,
             radius: bool = False):
    """Add a rectangle (or rounded rectangle) and return the shape."""
    if radius:
        # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE (5) for add_shape; arg 1 is rectangle
        sp = slide.shapes.add_shape(5, l, t, w, h)  # 5 = rounded rectangle preset
    else:
        sp = slide.shapes.add_shape(1, l, t, w, h)
    sp.text = ""
    if bg is not None:
        _set_solid_fill(sp, bg)
    else:
        sp.fill.background()
    if border is not None:
        _set_line(sp, border, border_pt)
    else:
        sp.line.fill.background()
    return sp


def _run(para, text, bold=False, italic=False, size_pt=None,
         color: RGBColor = None, font_name: str = None):
    """Append a run to a paragraph."""
    run = para.add_run()
    run.text = text
    rf = run.font
    if bold:     rf.bold   = True
    if italic:   rf.italic = True
    if size_pt:  rf.size   = Pt(size_pt)
    if color:    rf.color.rgb = color
    if font_name: rf.name  = font_name
    return run


def add_textbox(slide, l, t, w, h, word_wrap=True):
    """Add an empty text box and return (shape, tf)."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    return txb, tf


def _para(tf, text="", bold=False, italic=False, size_pt=None,
          color: RGBColor = None, align=PP_ALIGN.LEFT,
          font_name: str = "Calibri", space_before_pt=0, space_after_pt=0):
    """Add a paragraph to tf and return it."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    if space_before_pt:
        _set_para_spacing(p, before=space_before_pt, after=space_after_pt)
    _run(p, text, bold=bold, italic=italic, size_pt=size_pt,
         color=color, font_name=font_name)
    return p


def _set_para_spacing(para, before=0, after=0):
    """Set paragraph spacing (points)."""
    pPr = para._pPr
    if pPr is None:
        pPr = para._p.get_or_add_pPr()
    if before:
        spc = etree.SubElement(pPr, qn("a:spcBef"))
        spcPts = etree.SubElement(spc, qn("a:spcPts"))
        spcPts.set("val", str(int(before * 100)))
    if after:
        spc = etree.SubElement(pPr, qn("a:spcAft"))
        spcPts = etree.SubElement(spc, qn("a:spcPts"))
        spcPts.set("val", str(int(after * 100)))


# ─────────────────────────────────────────────────────────────────────────────
# Mid-level section builders (return next Y position in Emu)
# ─────────────────────────────────────────────────────────────────────────────

SECTION_HEAD_PT  = 34
BODY_PT          = 20
BULLET_PT        = 19
SMALL_PT         = 16
ROLE_TITLE_PT    = 17
ROLE_NAME_PT     = 20
SCENARIO_HEAD_PT = 20
SCENARIO_DESC_PT = 16
PILL_PT          = 16
METRIC_BIG_PT    = 44
METRIC_LBL_PT    = 15

GAP  = _i(0.08)   # small vertical gap between blocks
DIVH = _i(0.04)   # height of a thin divider line


def section_heading(slide, l, t, w, text: str) -> Emu:
    """Draw a centred red section heading; return next Y after."""
    h = _i(0.55)
    _, tf = add_textbox(slide, l, t, w, h)
    _para(tf, text, bold=True, size_pt=SECTION_HEAD_PT, color=RED,
          align=PP_ALIGN.CENTER, font_name="Calibri")
    return t + h + GAP


def divider(slide, l, t, w) -> Emu:
    """Draw a thin horizontal rule; return next Y after."""
    r = add_rect(slide, l, t + DIVH, w, DIVH, border=BORDER, border_pt=0.75)
    r.fill.background()
    return t + DIVH * 2 + GAP


def bullet_list(slide, l, t, w, items: list, title=None) -> Emu:
    """Draw a bullet list; return next Y after."""
    line_h = _i(0.30)
    h = line_h * (len(items) + (1 if title else 0)) + _i(0.1)
    _, tf = add_textbox(slide, l, t, w, h)
    tf.word_wrap = True
    first = True
    if title:
        p = _para(tf, title, bold=True, size_pt=BODY_PT, color=INK,
                  font_name="Calibri")
        first = False
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        # indent level 1 for bullet
        buNone = pPr.find(qn("a:buNone"))
        if buNone is not None:
            pPr.remove(buNone)
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "•")
        # indent: negative first-line indent for hanging bullet style
        pPr.set("indent", str(Emu(-_i(0.2))))
        pPr.set("marL",   str(Emu(_i(0.28))))
        _run(p, item, size_pt=BULLET_PT, color=INK, font_name="Calibri")
    return t + h + GAP


def team_members(slide, l, t, w) -> Emu:
    """Draw the 4 team member rows; return next Y after."""
    members = [
        ("TB", AV_RED,  "Backend & Real-Time Lead", "Tom Burchell",
         "Socket.IO engine · scenario logic · server architecture · CI/CD · cloud deployment"),
        ("LW", AV_BLUE, "Frontend & UI/UX Lead",    "Leon Wasiliew",
         "React UI · layouts · real-time state updates · responsive design"),
        ("JH", AV_TEAL, "Script Master / QA",       "John Hay",
         "Scenario writing · integration testing · accessibility review"),
        ("KP", AV_GOLD, "Project Manager",           "Kael Payette",
         "Coordination · documentation · Git workflow · standards"),
    ]
    avatar_sz = _i(0.55)
    row_h     = _i(0.90)
    gap_row   = _i(0.12)
    text_l    = l + avatar_sz + _i(0.12)
    text_w    = w - avatar_sz - _i(0.12)

    y = t
    for initials, color, role, name, tasks in members:
        # avatar circle (rounded rectangle approximation)
        av = slide.shapes.add_shape(5, l, y + _i(0.05),
                                     avatar_sz, avatar_sz)
        _set_solid_fill(av, color)
        av.line.fill.background()
        av.text_frame.word_wrap = True
        p = av.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, initials, bold=True, size_pt=18, color=WHITE,
             font_name="Calibri")
        av.text_frame.auto_size = None
        # auto-center vertically inside the avatar
        av.text_frame._txBody.bodyPr.set("anchor", "ctr")

        # role / name / tasks text
        _, tf = add_textbox(slide, text_l, y, text_w, row_h)
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        _run(p1, role, bold=True, size_pt=ROLE_TITLE_PT,
             color=RED_DARK, font_name="Calibri")
        p2 = tf.add_paragraph()
        _run(p2, name, bold=True, size_pt=ROLE_NAME_PT,
             color=INK, font_name="Calibri")
        p3 = tf.add_paragraph()
        _run(p3, tasks, size_pt=SMALL_PT, color=MUTED, font_name="Calibri")

        y += row_h + gap_row

    return y + GAP


def arch_diagram(slide, l, t, w) -> Emu:
    """Draw the architecture (ERD) diagram; return next Y after."""
    pad  = _i(0.15)
    il   = l + pad
    iw   = w - pad * 2
    y    = t + pad

    box_h  = _i(0.65)
    arrow_w= _i(0.55)

    # Row 1: Instructor   ⬦   Trainees
    box_w1 = (iw - arrow_w) / 2
    # Instructor box
    b1 = slide.shapes.add_shape(1, il, y, box_w1, box_h)
    _set_solid_fill(b1, AB_BLUE_BG); _set_line(b1, AB_BLUE_BD, 1.5)
    b1.text_frame.word_wrap = True
    b1.text_frame._txBody.bodyPr.set("anchor", "ctr")
    p = b1.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "Instructor Admin Dashboard", bold=True,
         size_pt=17, color=AB_BLUE_TX, font_name="Calibri")
    # diamond / connector arrow label
    _, tf_arr = add_textbox(slide, il + box_w1, y, arrow_w, box_h)
    _para(tf_arr, "⬦", size_pt=22, color=MUTED, align=PP_ALIGN.CENTER)
    tf_arr._txBody.bodyPr.set("anchor", "ctr")
    # Trainees box
    b2 = slide.shapes.add_shape(1, il + box_w1 + arrow_w, y, box_w1, box_h)
    _set_solid_fill(b2, AB_BLUE_BG); _set_line(b2, AB_BLUE_BD, 1.5)
    b2.text_frame.word_wrap = True
    b2.text_frame._txBody.bodyPr.set("anchor", "ctr")
    p = b2.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "Trainees React UI (×10+)", bold=True,
         size_pt=17, color=AB_BLUE_TX, font_name="Calibri")

    y += box_h + _i(0.05)
    # arrow label
    _, tf_mid = add_textbox(slide, il, y, iw, _i(0.28))
    _para(tf_mid, "↕ WebSocket / HTTP", size_pt=16, color=MUTED,
          align=PP_ALIGN.CENTER)
    y += _i(0.28) + _i(0.04)

    # Row 2: Socket.IO Server full-width
    b3 = slide.shapes.add_shape(1, il, y, iw, box_h)
    _set_solid_fill(b3, AB_GREEN_BG); _set_line(b3, AB_GREEN_BD, 1.5)
    b3.text_frame.word_wrap = True
    b3.text_frame._txBody.bodyPr.set("anchor", "ctr")
    p = b3.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "⚡ Socket.IO Server  ·  Express REST API  ·  Node.js 20",
         bold=True, size_pt=17, color=AB_GREEN_TX, font_name="Calibri")

    y += box_h + _i(0.05)
    _, tf_dn = add_textbox(slide, il, y, iw, _i(0.28))
    _para(tf_dn, "↕", size_pt=16, color=MUTED, align=PP_ALIGN.CENTER)
    y += _i(0.28) + _i(0.04)

    # Row 3: three sub-system boxes
    sub_w = (iw - _i(0.12)) / 3
    pairs = [
        ("In-Memory\nSession Store", AB_RED_BG, AB_RED_BD, AB_RED_TX),
        ("Escalation\nEngine",       AB_GOLD_BG, AB_GOLD_BD, AB_GOLD_TX),
        ("Scenario\nRegistry (8)",   AB_PURP_BG, AB_PURP_BD, AB_PURP_TX),
    ]
    for idx, (txt, bg, bd, tx) in enumerate(pairs):
        bx = slide.shapes.add_shape(
            1, il + idx * (sub_w + _i(0.06)), y, sub_w, box_h)
        _set_solid_fill(bx, bg); _set_line(bx, bd, 1.5)
        bx.text_frame.word_wrap = True
        bx.text_frame._txBody.bodyPr.set("anchor", "ctr")
        p = bx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, txt, bold=True, size_pt=16, color=tx, font_name="Calibri")

    y += box_h + _i(0.1)
    _, tf_note = add_textbox(slide, il, y, iw, _i(0.28))
    _para(tf_note,
          "Rooms isolate sessions  ·  Events broadcast to all participants"
          "  ·  Automatic timeline progression",
          italic=True, size_pt=14, color=MUTED, align=PP_ALIGN.CENTER)
    y += _i(0.28)

    # outer border
    outer_h = y - t + pad
    border_box = add_rect(slide, l, t, w, outer_h,
                          bg=RGBColor(0xE8, 0xF8, 0xEF),
                          border=AB_GREEN_BD, border_pt=1.5)
    # push border behind content just added
    sp_tree = slide.shapes._spTree
    border_elem = border_box._element
    sp_tree.remove(border_elem)
    sp_tree.insert(2, border_elem)   # index 2 keeps it behind slide content

    return y + pad + GAP


def scenario_grid(slide, l, t, w) -> Emu:
    """Draw the 8 scenario mini-cards (2-col grid); return next Y after."""
    scenarios = [
        ("Active Threat",     "Lockdown, sweep, all-clear procedures"),
        ("Fire Emergency",    "Alarm, evacuation, fire dept. handover"),
        ("Flood",             "Water intrusion, system protection"),
        ("Cyber Attack",      "Ransomware, network isolation, IR"),
        ("Power Outage",      "UPS/generator failover, critical systems"),
        ("Severe Weather",    "Shelter-in-place, facility protection"),
        ("Medical Emergency", "Cardiac event, first-aid, EMS coord."),
        ("Hazmat Spill",      "Containment, PPE, decontamination"),
    ]
    cols  = 2
    gap_x = _i(0.10)
    gap_y = _i(0.08)
    card_w = (w - gap_x) / cols
    card_h = _i(0.72)
    SC_BG  = RGBColor(0xFF, 0xF8, 0xF8)
    SC_BD  = RGBColor(0xF4, 0xC4, 0xC4)

    y = t
    for i, (name, desc) in enumerate(scenarios):
        col = i % cols
        row = i // cols
        cx  = l + col * (card_w + gap_x)
        cy  = y + row * (card_h + gap_y)

        bx = add_rect(slide, cx, cy, card_w, card_h, bg=SC_BG, border=SC_BD)
        _, tf = add_textbox(slide, cx + _i(0.1), cy + _i(0.06),
                            card_w - _i(0.2), card_h - _i(0.12))
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        _run(p1, name, bold=True, size_pt=SCENARIO_HEAD_PT,
             color=RED_DARK, font_name="Calibri")
        p2 = tf.add_paragraph()
        _run(p2, desc, size_pt=SCENARIO_DESC_PT, color=MUTED,
             font_name="Calibri")

    rows_used = (len(scenarios) + cols - 1) // cols
    total_h   = rows_used * (card_h + gap_y)
    return t + total_h + GAP


def app_mockups(slide, l, t, w) -> Emu:
    """Draw simple application mockup thumbnails; return next Y after."""
    labels = [
        ("Home Page",        RED,          "S.I.R.E. landing · scenario overview"),
        ("Login Page",       AV_BLUE,      "Email / password · role select"),
        ("Admin Dashboard",  RED,          "Session key · inject controls"),
        ("Trainee Interface",AV_TEAL,      "Live events · action submission"),
        ("Create Session",   AV_BLUE,      "Scenario select · start session"),
        ("Join Session",     AV_GOLD,      "6-char code · trainee onboarding"),
    ]
    cols   = 3
    gap_x  = _i(0.10)
    gap_y  = _i(0.10)
    card_w = (w - gap_x * (cols - 1)) / cols
    card_h = _i(0.85)
    CHROME = RGBColor(0xF8, 0xF8, 0xF8)
    BAR_H  = _i(0.12)

    y = t
    for i, (label, bar_color, sub) in enumerate(labels):
        col = i % cols
        row = i // cols
        cx  = l + col * (card_w + gap_x)
        cy  = y + row * (card_h + gap_y)

        # card frame
        add_rect(slide, cx, cy, card_w, card_h,
                 bg=CHROME, border=BORDER, border_pt=0.75)
        # browser bar stripe
        bar = add_rect(slide, cx, cy, card_w, BAR_H, bg=bar_color)
        bar.line.fill.background()
        # label
        _, tf = add_textbox(slide, cx + _i(0.05), cy + BAR_H + _i(0.04),
                            card_w - _i(0.10), card_h - BAR_H - _i(0.08))
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        _run(p1, label, bold=True, size_pt=15, color=INK, font_name="Calibri")
        p2 = tf.add_paragraph()
        _run(p2, sub, size_pt=13, color=MUTED, font_name="Calibri")

    rows_used = (len(labels) + cols - 1) // cols
    total_h   = rows_used * (card_h + gap_y)
    return t + total_h + GAP


def swot_grid(slide, l, t, w) -> Emu:
    """Draw the SWOT 2×2 grid; return next Y after."""
    items = [
        ("Strengths",
         "Real-time WebSocket; full-stack JS; modular scenarios;"
         " session isolation; no DB dependency",
         SWOT_S_BG, SWOT_S_BD, SWOT_S_TXT),
        ("Weaknesses",
         "In-memory storage (no restart persistence);"
         " no persistent accounts; WAN latency risk",
         SWOT_W_BG, SWOT_W_BD, SWOT_W_TXT),
        ("Opportunities",
         "Persistent session logs; mobile-responsive client;"
         " expanded scenario library; LMS integration",
         SWOT_O_BG, SWOT_O_BD, SWOT_O_TXT),
        ("Threats",
         "Free-tier cold-starts; browser WebSocket restrictions;"
         " server restart clears active sessions",
         SWOT_T_BG, SWOT_T_BD, SWOT_T_TXT),
    ]
    gap    = _i(0.08)
    cell_w = (w - gap) / 2
    cell_h = _i(1.30)

    for i, (title, text, bg, bd, tx) in enumerate(items):
        col = i % 2; row = i // 2
        cx  = l + col * (cell_w + gap)
        cy  = t + row * (cell_h + gap)
        add_rect(slide, cx, cy, cell_w, cell_h, bg=bg, border=bd)
        _, tf = add_textbox(slide, cx + _i(0.10), cy + _i(0.10),
                            cell_w - _i(0.20), cell_h - _i(0.20))
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        _run(p1, title, bold=True, size_pt=20, color=tx, font_name="Calibri")
        p2 = tf.add_paragraph()
        _run(p2, text, size_pt=SMALL_PT, color=tx, font_name="Calibri")

    rows_used = 2
    return t + rows_used * (cell_h + gap) + GAP


def tech_stack(slide, l, t, w) -> Emu:
    """Draw technology stack section; return next Y after."""
    sections = [
        ("FRONTEND",
         ["React 19", "React Router 7", "Vite 7", "Socket.IO Client 4"],
         FE_BG, FE_TX),
        ("BACKEND",
         ["Node.js 20", "Express 4", "Socket.IO 4", "ES Modules", "Nanoid"],
         BE_BG, BE_TX),
        ("REAL-TIME",
         ["WebSocket", "Long-Poll Fallback", "Room Isolation", "Auto-Reconnect"],
         RT_BG, RT_TX),
        ("DEVOPS & DEPLOYMENT",
         ["Docker (Alpine)", "GitHub Actions CI", "Railway", "Render", "Heroku"],
         OPS_BG, OPS_TX),
    ]
    y = t
    for label, pills, bg, tx in sections:
        # Label
        _, tf_lbl = add_textbox(slide, l, y, w, _i(0.28))
        _para(tf_lbl, label, bold=True, size_pt=14, color=MUTED,
              font_name="Calibri")
        y += _i(0.28)

        # Pills row — draw each pill as a small rect with text
        PILL_H = _i(0.32)
        PILL_PAD_X = _i(0.12)
        px = l
        for pill_text in pills:
            # approximate width: 1 char ≈ 0.10" at pill_pt
            pill_w = max(_i(0.8), len(pill_text) * _i(0.095) + PILL_PAD_X * 2)
            if px + pill_w > l + w:
                px  = l
                y  += PILL_H + _i(0.05)
            add_rect(slide, px, y, pill_w, PILL_H, bg=bg, border=tx,
                     border_pt=0.5)
            _, tf_p = add_textbox(slide, px, y, pill_w, PILL_H)
            tf_p.word_wrap = False
            tf_p._txBody.bodyPr.set("anchor", "ctr")
            p = tf_p.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            _run(p, pill_text, bold=True, size_pt=PILL_PT,
                 color=tx, font_name="Calibri")
            px += pill_w + _i(0.06)
        y += PILL_H + _i(0.10)

    return y + GAP


def key_metrics(slide, l, t, w) -> Emu:
    """Draw the 4 key-metric cards; return next Y after."""
    metrics = [
        ("8",      "SCENARIOS",      PINK_BG,   MET_RED),
        ("10+",    "TRAINEES/SESSION", BLUE_M_BG, MET_BLUE),
        ("<300ms", "EVENT LATENCY",  GRN_M_BG,  MET_GREEN),
        ("12",     "WEEKS DEVELOPMENT", YEL_M_BG,  MET_GOLD),
    ]
    gap    = _i(0.10)
    card_w = (w - gap) / 2
    card_h = _i(1.05)
    BORDER_C = RGBColor(0xDD, 0xDD, 0xDD)

    for i, (value, label, bg, val_color) in enumerate(metrics):
        col = i % 2; row = i // 2
        cx  = l + col * (card_w + gap)
        cy  = t + row * (card_h + gap)
        add_rect(slide, cx, cy, card_w, card_h, bg=bg, border=BORDER_C)
        # value
        _, tf_v = add_textbox(slide, cx + _i(0.05), cy + _i(0.12),
                              card_w - _i(0.10), _i(0.55))
        tf_v._txBody.bodyPr.set("anchor", "ctr")
        p = tf_v.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, value, bold=True, size_pt=METRIC_BIG_PT,
             color=val_color, font_name="Calibri")
        # label
        _, tf_l = add_textbox(slide, cx + _i(0.05), cy + _i(0.65),
                              card_w - _i(0.10), _i(0.32))
        tf_l._txBody.bodyPr.set("anchor", "ctr")
        p = tf_l.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, label, size_pt=METRIC_LBL_PT,
             color=MUTED, font_name="Calibri")

    rows_used = 2
    return t + rows_used * (card_h + gap) + GAP


def sire_logo_block(slide, l, t, w) -> Emu:
    """Draw a styled S.I.R.E. logo text block; return next Y after."""
    h = _i(1.50)
    bg_rect = add_rect(slide, l, t, w, h, bg=INK)
    bg_rect.line.fill.background()
    # Large S.I.R.E. title
    _, tf = add_textbox(slide, l + _i(0.15), t + _i(0.18),
                        w - _i(0.30), h - _i(0.36))
    tf.word_wrap = False
    tf._txBody.bodyPr.set("anchor", "ctr")
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    _run(p1, "S.I.R.E.", bold=True, size_pt=64,
         color=RED, font_name="Calibri")
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    _run(p2, "Simulated Incident Response Environment",
         size_pt=14, color=WHITE, font_name="Calibri")
    return t + h + GAP


# ─────────────────────────────────────────────────────────────────────────────
# Main assembly
# ─────────────────────────────────────────────────────────────────────────────

def build_presentation():
    prs  = Presentation(TEMPLATE)
    slide = prs.slides[0]

    # ── 1. Re-colour header bar to poster red ─────────────────────────────────
    for shape in slide.shapes:
        if shape.name == "Rounded Rectangle 53":
            _set_solid_fill(shape, RED)
            break

    # ── 2. Remove placeholder text boxes, old labels, and attribution box ──────
    remove_names = {
        "TextBox 57", "TextBox 58", "TextBox 59",
        "Rectangle 34", "Rectangle 35", "Rectangle 36",
        "Rectangle 54", "Rectangle 56",
    }
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.name in remove_names:
            sp_tree.remove(shape._element)

    # ── 3. Panel content boundaries (match template rounded-rect positions) ───
    # Left panel:   l=0.41", t=4.67", w=11.00", h=26.50"
    # Center panel: l=12.12", t=4.68", w=23.00", h=26.50"
    # Right panel:  l=35.86", t=4.47", w=11.00", h=26.50"
    PAD  = _i(0.30)   # inner horizontal padding

    LP_L = _i(0.41)  + PAD;  LP_T = _i(4.77);  LP_W = _i(11.00) - PAD * 2
    CP_L = _i(12.12) + PAD;  CP_T = _i(4.77);  CP_W = _i(23.00) - PAD * 2
    RP_L = _i(35.86) + PAD;  RP_T = _i(4.77);  RP_W = _i(11.00) - PAD * 2

    # ── 4. Header text ────────────────────────────────────────────────────────
    HDR_L = _i(0.60); HDR_T = _i(0.52); HDR_W = _i(43.00); HDR_H = _i(2.70)
    _, tf_hdr = add_textbox(slide, HDR_L, HDR_T, HDR_W, HDR_H)
    tf_hdr.word_wrap = True
    tf_hdr._txBody.bodyPr.set("anchor", "ctr")
    p1 = tf_hdr.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    _run(p1, "S.I.R.E. — Simulated Incident Response Environment",
         bold=True, size_pt=72, color=WHITE, font_name="Calibri")
    p2 = tf_hdr.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    _run(p2, "Team Nibble  ·  NSCC  ·  PROG 3300 — Integrated Projects for Programming"
         "  ·  Instructor: Alfred Parkes",
         bold=True, size_pt=30, color=RGBColor(0xFF, 0xDD, 0xDD),
         font_name="Calibri")
    p3 = tf_hdr.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    _run(p3, "A real-time emergency response training platform",
         italic=True, size_pt=24, color=RGBColor(0xFF, 0xCC, 0xCC),
         font_name="Calibri")

    # ── 5. LEFT PANEL: Features · Scope · Roles ───────────────────────────────
    y = LP_T

    y = section_heading(slide, LP_L, y, LP_W, "Features")
    y = bullet_list(slide, LP_L, y, LP_W, [
        "Real-time event broadcasting via WebSocket",
        "Scenario loading with timed, auto-escalating injects",
        "Admin controls — start, pause, escalate, end",
        "Trainee decision and action submission",
        "Branching scenario logic with severity levels",
        "Session key generation (6-character codes)",
        "Live progress monitoring for instructors",
        "Feedback and post-exercise results display",
        "Colour-coded severity indicators (info/warning/critical)",
        "Multi-user rooms with automatic reconnection",
    ])

    y = divider(slide, LP_L, y, LP_W)
    y = section_heading(slide, LP_L, y, LP_W, "Scope")
    y = bullet_list(slide, LP_L, y, LP_W, [
        "Full-stack JS web application",
        "Admin dashboard + trainee interface",
        "Real-time WebSocket communication",
        "Eight complete emergency scenarios",
        "Session-based access (admin/trainee)",
        "Free-tier cloud deployment",
        "Documentation + testing artefacts",
    ])

    y = divider(slide, LP_L, y, LP_W)
    y = section_heading(slide, LP_L, y, LP_W, "Roles")
    team_members(slide, LP_L, y, LP_W)

    # ── 6. CENTRE PANEL: Description · Architecture · Scenarios · Mockups ─────
    y = CP_T

    y = section_heading(slide, CP_L, y, CP_W, "Description")
    _, tf_desc = add_textbox(slide, CP_L, y, CP_W, _i(1.40))
    tf_desc.word_wrap = True
    _para(tf_desc,
          "A full-stack, real-time simulation platform for practising emergency "
          "response scenarios. Built with Node.js, Express, Socket.IO, and React, "
          "SIRE enables instructors to run dynamic, escalating incident simulations "
          "while trainees respond in real time.",
          size_pt=BODY_PT, color=INK, align=PP_ALIGN.CENTER, font_name="Calibri")
    y += _i(1.40) + GAP

    y = divider(slide, CP_L, y, CP_W)
    y = section_heading(slide, CP_L, y, CP_W, "Architecture (ERD)")
    y = arch_diagram(slide, CP_L, y, CP_W)

    y = divider(slide, CP_L, y, CP_W)
    y = section_heading(slide, CP_L, y, CP_W, "Training Scenarios")
    y = scenario_grid(slide, CP_L, y, CP_W)

    y = divider(slide, CP_L, y, CP_W)
    y = section_heading(slide, CP_L, y, CP_W, "Application Mockups")
    app_mockups(slide, CP_L, y, CP_W)

    # ── 7. RIGHT PANEL: Logo · SWOT · Tech Stack · Key Metrics ───────────────
    y = RP_T

    y = sire_logo_block(slide, RP_L, y, RP_W)

    y = divider(slide, RP_L, y, RP_W)
    y = section_heading(slide, RP_L, y, RP_W, "SWOT")
    y = swot_grid(slide, RP_L, y, RP_W)

    y = divider(slide, RP_L, y, RP_W)
    y = section_heading(slide, RP_L, y, RP_W, "Technology Stack")
    y = tech_stack(slide, RP_L, y, RP_W)

    y = divider(slide, RP_L, y, RP_W)
    y = section_heading(slide, RP_L, y, RP_W, "Key Metrics")
    key_metrics(slide, RP_L, y, RP_W)

    # ── 8. Full-width footer bar (matches poster footer style) ───────────────
    # Template footer sits below the panels: ~32.6" from top
    FOOTER_T = _i(32.00)
    FOOTER_H = _i(2.20)
    FULL_W   = _i(47.00)   # matches header bar width

    # Full-width red background (same as header)
    ftr_bg = slide.shapes.add_shape(1, _i(0.12), FOOTER_T, FULL_W, FOOTER_H)
    _set_solid_fill(ftr_bg, RED)
    ftr_bg.line.fill.background()

    # Left section: project name + course credits
    LF_L = _i(0.50); LF_W = _i(30.00)
    _, tf_lf = add_textbox(slide, LF_L, FOOTER_T + _i(0.25),
                           LF_W, FOOTER_H - _i(0.40))
    tf_lf.word_wrap = True
    tf_lf._txBody.bodyPr.set("anchor", "ctr")
    p_lf1 = tf_lf.paragraphs[0]; p_lf1.alignment = PP_ALIGN.LEFT
    _run(p_lf1, "S.I.R.E. — Simulated Incident Response Environment",
         bold=True, size_pt=26, color=WHITE, font_name="Calibri")
    p_lf2 = tf_lf.add_paragraph(); p_lf2.alignment = PP_ALIGN.LEFT
    _run(p_lf2,
         "NSCC · PROG 3300 Integrated Projects for Programming · "
         "Instructor: Alfred Parkes  |  "
         "Team: Tom Burchell · Leon Wasiliew · John Hay · Kael Payette",
         size_pt=20, color=RGBColor(0xFF, 0xDD, 0xDD), font_name="Calibri")
    p_lf3 = tf_lf.add_paragraph(); p_lf3.alignment = PP_ALIGN.LEFT
    _run(p_lf3, "Scenario images: Gemini  ·  github.com/Anaxagorius/S.I.R.E.",
         italic=True, size_pt=16,
         color=RGBColor(0xFF, 0xCC, 0xCC), font_name="Calibri")

    # Right section: tech badges
    badge_items = ["Node.js 20", "React 19", "Socket.IO 4", "Docker"]
    BADGE_H   = _i(0.55); BADGE_GAP = _i(0.10)
    # measure total badge width to right-align
    badge_widths = [max(_i(1.10), len(b) * _i(0.115) + _i(0.25)) for b in badge_items]
    total_bw  = sum(badge_widths) + BADGE_GAP * (len(badge_items) - 1)
    bx = _i(0.12) + FULL_W - total_bw - _i(0.40)
    by = FOOTER_T + (FOOTER_H - BADGE_H) / 2

    for badge, bw in zip(badge_items, badge_widths):
        br = add_rect(slide, bx, by, bw, BADGE_H,
                      bg=RGBColor(0x33, 0x33, 0x33),
                      border=RGBColor(0x88, 0x88, 0x88), border_pt=0.5)
        _, tf_b = add_textbox(slide, bx, by, bw, BADGE_H)
        tf_b.word_wrap = False
        tf_b._txBody.bodyPr.set("anchor", "ctr")
        p_b = tf_b.paragraphs[0]; p_b.alignment = PP_ALIGN.CENTER
        _run(p_b, badge, bold=True, size_pt=18,
             color=WHITE, font_name="Calibri")
        bx += bw + BADGE_GAP

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(OUTPUT)
    size_kb = os.path.getsize(OUTPUT) // 1024
    print(f"✅  Presentation saved → {OUTPUT}  ({size_kb} KB)")
    print("    Open in Microsoft PowerPoint or LibreOffice Impress.")


if __name__ == "__main__":
    build_presentation()
